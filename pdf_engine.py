"""Pure-Python PDF conversions used by the pywebview API bridge."""

import io
import os


def _safe_stem(path):
    return os.path.splitext(os.path.basename(path))[0]


def render_pdf_pages(source, output_dir, image_format="PNG", dpi=150):
    import pypdfium2 as pdfium

    document = pdfium.PdfDocument(source)
    extension = ".jpg" if image_format == "JPEG" else ".png"
    outputs = []
    try:
        for number in range(len(document)):
            bitmap = document[number].render(scale=max(72, dpi) / 72)
            image = bitmap.to_pil()
            if image_format == "JPEG" and image.mode != "RGB":
                image = image.convert("RGB")
            destination = os.path.join(output_dir, f"{_safe_stem(source)}-page-{number + 1}{extension}")
            save_options = {"quality": 92} if image_format == "JPEG" else {}
            image.save(destination, image_format, **save_options)
            outputs.append(destination)
    finally:
        document.close()
    return outputs


def extract_pdf_text(source, destination):
    import pdfplumber

    with pdfplumber.open(source) as document:
        text = "\n\n".join(page.extract_text() or "" for page in document.pages)
    with open(destination, "w", encoding="utf-8") as output:
        output.write(text)
    return text


def merge_pdfs(sources, destination):
    from pypdf import PdfWriter

    writer = PdfWriter()
    try:
        for source in sources:
            writer.append(source)
        with open(destination, "wb") as output:
            writer.write(output)
    finally:
        writer.close()


def split_pdf(source, output_dir):
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(source)
    outputs = []
    for index, page in enumerate(reader.pages, start=1):
        destination = os.path.join(output_dir, f"{_safe_stem(source)}-page-{index}.pdf")
        writer = PdfWriter()
        try:
            writer.add_page(page)
            with open(destination, "wb") as output:
                writer.write(output)
        finally:
            writer.close()
        outputs.append(destination)
    return outputs


def rotate_pdf(source, destination, degrees):
    from pypdf import PdfReader, PdfWriter

    reader, writer = PdfReader(source), PdfWriter()
    try:
        for page in reader.pages:
            writer.add_page(page.rotate(degrees))
        with open(destination, "wb") as output:
            writer.write(output)
    finally:
        writer.close()


def extract_pdf_pages(source, destination, pages):
    from pypdf import PdfReader, PdfWriter

    reader, writer = PdfReader(source), PdfWriter()
    try:
        for page_number in pages:
            if page_number < 1 or page_number > len(reader.pages):
                raise ValueError(f"页码 {page_number} 超出范围（共 {len(reader.pages)} 页）")
            writer.add_page(reader.pages[page_number - 1])
        with open(destination, "wb") as output:
            writer.write(output)
    finally:
        writer.close()


def images_to_pdf(sources, destination):
    from PIL import Image

    images = []
    try:
        for source in sources:
            image = Image.open(source).convert("RGB")
            images.append(image.copy())
            image.close()
        if not images:
            raise ValueError("请至少选择一张图片")
        images[0].save(destination, "PDF", save_all=True, append_images=images[1:], resolution=150)
    finally:
        for image in images:
            image.close()


def pdf_to_docx(source, destination):
    import pdfplumber
    from docx import Document

    document = Document()
    wrote_text = False
    with pdfplumber.open(source) as pdf:
        for index, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if text:
                wrote_text = True
                for line in text.splitlines():
                    document.add_paragraph(line)
            if index < len(pdf.pages) - 1:
                document.add_page_break()
    if not wrote_text:
        raise ValueError("未提取到可编辑文字；该 PDF 可能是扫描件或纯图片")
    document.save(destination)


def pdf_to_xlsx(source, destination):
    import pdfplumber
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    table_count = 0
    with pdfplumber.open(source) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            for table_number, table in enumerate(page.extract_tables(), start=1):
                if not table or not any(any(cell for cell in row) for row in table):
                    continue
                sheet = workbook.create_sheet(f"第{page_number}页表{table_number}")
                for row in table:
                    sheet.append([cell or "" for cell in row])
                table_count += 1
    if not table_count:
        raise ValueError("未发现规整表格，未生成空白 Excel")
    workbook.save(destination)


def pdf_to_pptx(source, destination, dpi=150):
    import pypdfium2 as pdfium
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    document = pdfium.PdfDocument(source)
    try:
        for page in document:
            image = page.render(scale=max(72, dpi) / 72).to_pil().convert("RGB")
            stream = io.BytesIO()
            image.save(stream, "PNG")
            slide = presentation.slides.add_slide(blank)
            slide.shapes.add_picture(stream, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    finally:
        document.close()
    presentation.save(destination)
